"""
Interactive PPTX Split & Merge CLI Tool
Complete workflow: Split presentation → Modify → Merge back
With file path input and formatting preservation

Installation:
    pip install python-pptx

Usage:
    python pptx_cli_tool.py
"""

from pptx import Presentation
from copy import deepcopy
import os
from pathlib import Path
from typing import List
import logging
import sys


class PPTXSplitMerge:
    """Robust PPTX file splitter and merger with full formatting preservation."""
    
    def __init__(self):
        self.logger = self._setup_logger()
    
    def _setup_logger(self):
        """Setup logging for operations."""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s'
        )
        return logging.getLogger(__name__)
    
    def split_pptx(self, input_file: str, output_dir: str) -> List[str]:
        """
        Split a PPTX file - one slide per file.
        
        Args:
            input_file: Path to input PPTX file
            output_dir: Directory to save split files
        
        Returns:
            List of output file paths
        """
        try:
            # Validate input file
            if not os.path.exists(input_file):
                raise FileNotFoundError(f"File not found: {input_file}")
            
            if not input_file.lower().endswith('.pptx'):
                raise ValueError("File must be a .pptx file")
            
            # Load presentation
            prs = Presentation(input_file)
            total_slides = len(prs.slides)
            
            print(f"\n✓ Loading: {input_file}")
            print(f"  Total slides: {total_slides}")
            
            # Create output directory
            Path(output_dir).mkdir(parents=True, exist_ok=True)
            
            output_files = []
            basename = Path(input_file).stem
            
            # Split each slide into separate file
            for slide_idx, slide in enumerate(prs.slides, 1):
                # Create new presentation
                new_prs = Presentation()
                new_prs.slide_width = prs.slide_width
                new_prs.slide_height = prs.slide_height
                
                # Copy master slides
                self._copy_master_slides(prs, new_prs)
                
                # Clone this slide
                self._clone_slide(slide, new_prs)
                
                # Save output file with slide number
                output_filename = os.path.join(
                    output_dir,
                    f"{basename}_slide_{slide_idx:03d}.pptx"
                )
                new_prs.save(output_filename)
                output_files.append(output_filename)
                
                print(f"  ✓ Created: {os.path.basename(output_filename)}")
            
            print(f"\n✓ Split complete!")
            print(f"  Output directory: {os.path.abspath(output_dir)}")
            print(f"  Total files created: {len(output_files)}")
            
            return output_files
            
        except Exception as e:
            print(f"\n✗ Error splitting PPTX: {str(e)}")
            self.logger.error(f"Error splitting PPTX: {str(e)}")
            raise
    
    def merge_pptx(self, input_files: List[str], output_file: str) -> str:
        """
        Merge multiple PPTX files into one.
        
        Args:
            input_files: List of PPTX file paths to merge (in order)
            output_file: Output PPTX file path
        
        Returns:
            Path to output file
        """
        try:
            if not input_files:
                raise ValueError("No input files provided")
            
            # Validate all input files exist
            for input_file in input_files:
                if not os.path.exists(input_file):
                    raise FileNotFoundError(f"File not found: {input_file}")
                if not input_file.lower().endswith('.pptx'):
                    raise ValueError(f"File must be .pptx: {input_file}")
            
            print(f"\n✓ Starting merge...")
            print(f"  Files to merge: {len(input_files)}")
            
            # Load first presentation as base
            base_prs = Presentation(input_files[0])
            print(f"  ✓ Base: {os.path.basename(input_files[0])}")
            
            # Merge remaining presentations
            for input_file in input_files[1:]:
                source_prs = Presentation(input_file)
                print(f"  + Merging: {os.path.basename(input_file)}")
                
                for slide in source_prs.slides:
                    self._clone_slide(slide, base_prs)
            
            # Create output directory if needed
            Path(output_file).parent.mkdir(parents=True, exist_ok=True)
            
            # Save merged presentation
            base_prs.save(output_file)
            
            print(f"\n✓ Merge complete!")
            print(f"  Output file: {os.path.abspath(output_file)}")
            print(f"  Total slides: {len(base_prs.slides)}")
            
            return output_file
            
        except Exception as e:
            print(f"\n✗ Error merging PPTX: {str(e)}")
            self.logger.error(f"Error merging PPTX: {str(e)}")
            raise
    
    # ==================== HELPER METHODS ====================
    
    def _copy_master_slides(self, source_prs: Presentation, 
                           dest_prs: Presentation) -> None:
        """Copy master slides and layouts from source to destination."""
        try:
            for master in source_prs.slide_masters:
                for layout in master.slide_layouts:
                    dest_prs.slide_layouts.get_or_add_layout(layout)
        except Exception as e:
            self.logger.warning(f"Could not fully copy masters: {str(e)}")
    
    def _clone_slide(self, source_slide, dest_prs: Presentation) -> None:
        """Clone a slide with all formatting, shapes, and properties."""
        try:
            slide_layout = source_slide.slide_layout
            new_slide = dest_prs.slides.add_slide(slide_layout)
            
            # Copy slide-level properties
            self._copy_slide_properties(source_slide, new_slide)
            
            # Clear default placeholders
            for shape in list(new_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)
            
            # Clone all shapes from source
            for shape in source_slide.shapes:
                self._clone_shape(shape, new_slide)
            
        except Exception as e:
            self.logger.error(f"Error cloning slide: {str(e)}")
            raise
    
    def _clone_shape(self, source_shape, dest_slide) -> None:
        """Clone a shape preserving all formatting."""
        try:
            # Deep copy shape element (preserves all XML properties)
            shape_element = deepcopy(source_shape.element)
            dest_slide.shapes._spTree.insert_element_before(
                shape_element, 'p:extLst'
            )
            
        except Exception as e:
            self.logger.warning(f"Could not clone shape, using fallback: {e}")
            self._clone_shape_fallback(source_shape, dest_slide)
    
    def _clone_shape_fallback(self, source_shape, dest_slide) -> None:
        """Fallback method for cloning shapes."""
        try:
            if not source_shape.has_text_frame:
                return
            
            left = source_shape.left
            top = source_shape.top
            width = source_shape.width
            height = source_shape.height
            
            text_box = dest_slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            
            # Copy text frame properties
            text_frame.word_wrap = source_shape.text_frame.word_wrap
            text_frame.margin_bottom = source_shape.text_frame.margin_bottom
            text_frame.margin_top = source_shape.text_frame.margin_top
            text_frame.margin_left = source_shape.text_frame.margin_left
            text_frame.margin_right = source_shape.text_frame.margin_right
            text_frame.vertical_anchor = source_shape.text_frame.vertical_anchor
            
            # Copy paragraphs with formatting
            text_frame.clear()
            for source_para in source_shape.text_frame.paragraphs:
                self._copy_paragraph(source_para, text_frame)
            
            # Copy shape fill and line
            self._copy_shape_fill(source_shape, text_box)
            self._copy_shape_line(source_shape, text_box)
            
        except Exception as e:
            self.logger.warning(f"Fallback clone also failed: {str(e)}")
    
    def _copy_paragraph(self, source_para, dest_text_frame) -> None:
        """Copy paragraph with all text runs and formatting."""
        try:
            dest_para = dest_text_frame.add_paragraph()
            dest_para.text = ""
            
            # Copy paragraph properties
            dest_para.level = source_para.level
            dest_para.alignment = source_para.alignment
            dest_para.space_before = source_para.space_before
            dest_para.space_after = source_para.space_after
            dest_para.line_spacing = source_para.line_spacing
            
            # Copy runs with formatting
            dest_para.clear()
            for source_run in source_para.runs:
                self._copy_run(source_run, dest_para)
        
        except Exception as e:
            self.logger.warning(f"Error copying paragraph: {str(e)}")
    
    def _copy_run(self, source_run, dest_para) -> None:
        """Copy text run with character-level formatting."""
        try:
            dest_run = dest_para.add_run()
            dest_run.text = source_run.text
            
            # Font properties
            font_src = source_run.font
            font_dst = dest_run.font
            
            font_dst.name = font_src.name
            font_dst.size = font_src.size
            font_dst.bold = font_src.bold
            font_dst.italic = font_src.italic
            font_dst.underline = font_src.underline
            
            # Color
            if font_src.color.type:
                if hasattr(font_src.color, 'rgb'):
                    font_dst.color.rgb = font_src.color.rgb
            
            # Other properties
            if hasattr(font_src, 'all_caps'):
                font_dst.all_caps = font_src.all_caps
            if hasattr(font_src, 'subscript'):
                font_dst.subscript = font_src.subscript
            if hasattr(font_src, 'superscript'):
                font_dst.superscript = font_src.superscript
                
        except Exception as e:
            self.logger.warning(f"Error copying run: {str(e)}")
    
    def _copy_shape_fill(self, source_shape, dest_shape) -> None:
        """Copy shape fill properties."""
        try:
            if hasattr(source_shape, 'fill'):
                src_fill = source_shape.fill
                dst_fill = dest_shape.fill
                
                if hasattr(src_fill, 'type'):
                    dst_fill.solid()
                    if hasattr(src_fill, 'fore_color'):
                        try:
                            dst_fill.fore_color.rgb = src_fill.fore_color.rgb
                        except:
                            pass
        except Exception as e:
            self.logger.debug(f"Could not copy fill: {str(e)}")
    
    def _copy_shape_line(self, source_shape, dest_shape) -> None:
        """Copy shape outline/line properties."""
        try:
            if hasattr(source_shape, 'line'):
                src_line = source_shape.line
                dst_line = dest_shape.line
                
                if hasattr(src_line, 'color'):
                    try:
                        dst_line.color.rgb = src_line.color.rgb
                    except:
                        pass
                if hasattr(src_line, 'width'):
                    dst_line.width = src_line.width
        except Exception as e:
            self.logger.debug(f"Could not copy line: {str(e)}")
    
    def _copy_slide_properties(self, source_slide, dest_slide) -> None:
        """Copy slide-level properties like background."""
        try:
            if hasattr(source_slide.background, 'fill'):
                src_fill = source_slide.background.fill
                dst_fill = dest_slide.background.fill
                
                if hasattr(src_fill, 'solid'):
                    try:
                        dst_fill.solid()
                        dst_fill.fore_color.rgb = src_fill.fore_color.rgb
                    except:
                        pass
        except Exception as e:
            self.logger.debug(f"Could not copy slide properties: {str(e)}")


class PPTXCliTool:
    """Interactive CLI tool for PPTX split/merge operations."""
    
    def __init__(self):
        self.processor = PPTXSplitMerge()
    
    def display_banner(self):
        """Display welcome banner."""
        print("\n" + "="*60)
        print("  PPTX SPLIT & MERGE TOOL")
        print("  Split presentations and merge them back")
        print("="*60)
    
    def display_menu(self):
        """Display main menu options."""
        print("\nOptions:")
        print("  1. Split a PPTX file (one slide per file)")
        print("  2. Merge split PPTX files back together")
        print("  3. Exit")
        print("-" * 60)
    
    def get_file_path(self, prompt: str) -> str:
        """Get and validate file path from user."""
        while True:
            path = input(prompt).strip()
            
            if not path:
                print("✗ Path cannot be empty")
                continue
            
            # Remove quotes if user added them
            path = path.strip('"\'')
            
            if not os.path.exists(path):
                print(f"✗ File not found: {path}")
                continue
            
            return path
    
    def get_directory_path(self, prompt: str) -> str:
        """Get directory path from user."""
        while True:
            path = input(prompt).strip()
            
            if not path:
                print("✗ Path cannot be empty")
                continue
            
            # Remove quotes if user added them
            path = path.strip('"\'')
            
            return path
    
    def split_workflow(self):
        """Handle split workflow."""
        print("\n" + "="*60)
        print("SPLIT WORKFLOW")
        print("="*60)
        
        # Get input file path
        input_file = self.get_file_path(
            "\nEnter path to PPTX file to split: "
        )
        
        # Get output directory
        output_dir = self.get_directory_path(
            "Enter output directory (default: ./split_slides): "
        ) or "./split_slides"
        
        try:
            output_files = self.processor.split_pptx(input_file, output_dir)
            self._save_file_list(output_files, output_dir)
            
        except Exception as e:
            print(f"\n✗ Split failed: {str(e)}")
    
    def merge_workflow(self):
        """Handle merge workflow."""
        print("\n" + "="*60)
        print("MERGE WORKFLOW")
        print("="*60)
        
        # Get method of input
        print("\nHow to provide split files?")
        print("  1. From a directory (automatic detection)")
        print("  2. Enter file paths manually")
        print("-" * 60)
        
        choice = input("Choose option (1-2): ").strip()
        
        if choice == "1":
            self._merge_from_directory()
        elif choice == "2":
            self._merge_from_manual_input()
        else:
            print("✗ Invalid choice")
    
    def _merge_from_directory(self):
        """Merge files from a directory."""
        split_dir = self.get_file_path(
            "\nEnter path to directory with split PPTX files: "
        )
        
        # Get all PPTX files from directory (sorted by name)
        pptx_files = sorted([
            os.path.join(split_dir, f)
            for f in os.listdir(split_dir)
            if f.endswith('.pptx')
        ])
        
        if not pptx_files:
            print(f"✗ No PPTX files found in: {split_dir}")
            return
        
        print(f"\n✓ Found {len(pptx_files)} PPTX files:")
        for i, file in enumerate(pptx_files, 1):
            print(f"  {i}. {os.path.basename(file)}")
        
        output_file = self.get_directory_path(
            "\nEnter output file path (default: ./merged.pptx): "
        ) or "./merged.pptx"
        
        try:
            self.processor.merge_pptx(pptx_files, output_file)
        except Exception as e:
            print(f"\n✗ Merge failed: {str(e)}")
    
    def _merge_from_manual_input(self):
        """Merge files by manual path entry."""
        input_files = []
        
        print("\nEnter paths to PPTX files (one per line, empty line to finish):")
        while True:
            file_path = input(f"File {len(input_files) + 1}: ").strip()
            
            if not file_path:
                if input_files:
                    break
                else:
                    print("✗ Please enter at least one file path")
                    continue
            
            # Remove quotes
            file_path = file_path.strip('"\'')
            
            if not os.path.exists(file_path):
                print(f"✗ File not found: {file_path}")
                continue
            
            if not file_path.lower().endswith('.pptx'):
                print("✗ File must be a .pptx file")
                continue
            
            input_files.append(file_path)
            print(f"  ✓ Added: {os.path.basename(file_path)}")
        
        if len(input_files) < 2:
            print("✗ Need at least 2 files to merge")
            return
        
        output_file = input("\nEnter output file path (default: ./merged.pptx): ").strip() or "./merged.pptx"
        
        try:
            self.processor.merge_pptx(input_files, output_file)
        except Exception as e:
            print(f"\n✗ Merge failed: {str(e)}")
    
    def _save_file_list(self, files: List[str], directory: str):
        """Save list of split files to a text file."""
        try:
            list_file = os.path.join(directory, "split_files_list.txt")
            with open(list_file, 'w') as f:
                f.write("Split PPTX Files List\n")
                f.write("=" * 50 + "\n\n")
                for file in files:
                    f.write(f"{file}\n")
            
            print(f"\n✓ File list saved to: {list_file}")
        except Exception as e:
            print(f"\n⚠ Could not save file list: {str(e)}")
    
    def run(self):
        """Run the interactive CLI tool."""
        self.display_banner()
        
        while True:
            self.display_menu()
            choice = input("Choose option (1-3): ").strip()
            
            if choice == "1":
                self.split_workflow()
            elif choice == "2":
                self.merge_workflow()
            elif choice == "3":
                print("\n✓ Thank you for using PPTX Split & Merge Tool!")
                print("="*60 + "\n")
                break
            else:
                print("✗ Invalid choice. Please enter 1, 2, or 3")


# ==================== MAIN ====================

if __name__ == "__main__":
    try:
        cli = PPTXCliTool()
        cli.run()
    except KeyboardInterrupt:
        print("\n\n✗ Operation cancelled by user")
        sys.exit(0)
    except Exception as e:
        print(f"\n✗ Fatal error: {str(e)}")
        sys.exit(1)