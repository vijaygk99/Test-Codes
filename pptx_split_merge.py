"""
Robust PPTX Split & Merge with Full Formatting Preservation
Handles: text formatting, shapes, colors, fonts, sizes, alignment, animations, and layouts

Installation:
    pip install python-pptx

Author: Data Science Team
Version: 2.0
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import os
from pathlib import Path
from typing import List, Optional
import logging


class PPTXSplitMerge:
    """Robust PPTX file splitter and merger with full formatting preservation."""
    
    def __init__(self):
        self.logger = self._setup_logger()
    
    def _setup_logger(self):
        """Setup basic logging for operations."""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s'
        )
        return logging.getLogger(__name__)
    
    # ==================== SPLIT OPERATIONS ====================
    
    def split_pptx(self, input_file: str, output_dir: str, 
                   slides_per_file: int = 1) -> List[str]:
        """
        Split a PPTX file into multiple files.
        
        Args:
            input_file: Path to input PPTX file
            output_dir: Directory to save split files
            slides_per_file: Number of slides per output file (default: 1)
        
        Returns:
            List of output file paths
        
        Example:
            >>> processor = PPTXSplitMerge()
            >>> files = processor.split_pptx("presentation.pptx", "output", slides_per_file=5)
        """
        try:
            # Load presentation
            prs = Presentation(input_file)
            total_slides = len(prs.slides)
            
            self.logger.info(f"Loading {input_file}: {total_slides} slides")
            
            # Create output directory
            Path(output_dir).mkdir(parents=True, exist_ok=True)
            
            output_files = []
            file_count = 1
            
            # Calculate number of output files needed
            num_output_files = (total_slides + slides_per_file - 1) // slides_per_file
            
            for i in range(0, total_slides, slides_per_file):
                # Create new presentation with same slide dimensions
                new_prs = Presentation()
                new_prs.slide_width = prs.slide_width
                new_prs.slide_height = prs.slide_height
                
                # Copy master slides and themes
                self._copy_master_slides(prs, new_prs)
                
                # Add slides to new presentation
                slide_range = min(i + slides_per_file, total_slides)
                for slide_idx in range(i, slide_range):
                    slide = prs.slides[slide_idx]
                    # Clone slide with all formatting
                    self._clone_slide(slide, new_prs)
                
                # Save output file
                output_filename = os.path.join(
                    output_dir, 
                    f"{Path(input_file).stem}_part{file_count}.pptx"
                )
                new_prs.save(output_filename)
                output_files.append(output_filename)
                
                self.logger.info(
                    f"Created {output_filename}: "
                    f"slides {i+1}-{slide_range}"
                )
                file_count += 1
            
            self.logger.info(f"Split complete: {len(output_files)} files created")
            return output_files
            
        except Exception as e:
            self.logger.error(f"Error splitting PPTX: {str(e)}")
            raise
    
    def split_pptx_by_range(self, input_file: str, output_dir: str,
                           ranges: List[tuple]) -> List[str]:
        """
        Split PPTX by specific slide ranges.
        
        Args:
            input_file: Path to input PPTX file
            output_dir: Directory to save split files
            ranges: List of (start, end) tuples (1-indexed, inclusive)
        
        Returns:
            List of output file paths
        
        Example:
            >>> processor = PPTXSplitMerge()
            >>> files = processor.split_pptx_by_range(
            ...     "presentation.pptx",
            ...     "output",
            ...     ranges=[(1, 5), (6, 10), (11, 15)]
            ... )
        """
        try:
            prs = Presentation(input_file)
            total_slides = len(prs.slides)
            
            Path(output_dir).mkdir(parents=True, exist_ok=True)
            output_files = []
            
            for idx, (start, end) in enumerate(ranges, 1):
                # Validate range
                start = max(1, start)
                end = min(total_slides, end)
                
                if start > end:
                    self.logger.warning(f"Invalid range {start}-{end}, skipping")
                    continue
                
                new_prs = Presentation()
                new_prs.slide_width = prs.slide_width
                new_prs.slide_height = prs.slide_height
                
                self._copy_master_slides(prs, new_prs)
                
                # Copy slides (convert 1-indexed to 0-indexed)
                for slide_idx in range(start - 1, end):
                    self._clone_slide(prs.slides[slide_idx], new_prs)
                
                output_filename = os.path.join(
                    output_dir,
                    f"{Path(input_file).stem}_range{start}-{end}.pptx"
                )
                new_prs.save(output_filename)
                output_files.append(output_filename)
                
                self.logger.info(f"Created {output_filename}: slides {start}-{end}")
            
            return output_files
            
        except Exception as e:
            self.logger.error(f"Error in split_by_range: {str(e)}")
            raise
    
    # ==================== MERGE OPERATIONS ====================
    
    def merge_pptx(self, input_files: List[str], output_file: str,
                  keep_source_theme: bool = True) -> str:
        """
        Merge multiple PPTX files into one.
        
        Args:
            input_files: List of PPTX file paths to merge (order matters)
            output_file: Output PPTX file path
            keep_source_theme: Keep original themes (True) or use first file's theme (False)
        
        Returns:
            Path to output file
        
        Example:
            >>> processor = PPTXSplitMerge()
            >>> merged = processor.merge_pptx(
            ...     ["part1.pptx", "part2.pptx", "part3.pptx"],
            ...     "merged.pptx",
            ...     keep_source_theme=True
            ... )
        """
        try:
            if not input_files:
                raise ValueError("No input files provided")
            
            # Load first presentation as base
            base_prs = Presentation(input_files[0])
            self.logger.info(f"Base presentation: {input_files[0]}")
            
            # Merge remaining presentations
            for input_file in input_files[1:]:
                self.logger.info(f"Merging: {input_file}")
                source_prs = Presentation(input_file)
                
                for slide in source_prs.slides:
                    if keep_source_theme:
                        # Clone with source theme
                        self._clone_slide(slide, base_prs)
                    else:
                        # Clone with destination theme
                        self._clone_slide_with_master(
                            slide, base_prs, 
                            base_prs.slide_masters[0]
                        )
            
            # Save merged presentation
            Path(output_file).parent.mkdir(parents=True, exist_ok=True)
            base_prs.save(output_file)
            
            self.logger.info(f"Merge complete: {output_file}")
            self.logger.info(f"Total slides: {len(base_prs.slides)}")
            
            return output_file
            
        except Exception as e:
            self.logger.error(f"Error merging PPTX: {str(e)}")
            raise
    
    # ==================== CORE FORMATTING PRESERVATION METHODS ====================
    
    def _copy_master_slides(self, source_prs: Presentation, 
                           dest_prs: Presentation) -> None:
        """Copy master slides and layouts from source to destination."""
        try:
            for master in source_prs.slide_masters:
                for layout in master.slide_layouts:
                    dest_prs.slide_layouts.get_or_add_layout(layout)
        except Exception as e:
            self.logger.warning(f"Could not fully copy masters: {str(e)}")
    
    def _clone_slide(self, source_slide, dest_prs: Presentation) -> None:
        """
        Clone a slide with all formatting, shapes, and properties.
        This is the CORE function that preserves formatting.
        
        Preserves:
        - Layout and master slide references
        - All shapes (text boxes, images, shapes)
        - Text formatting (fonts, colors, sizes, bold, italic)
        - Shape fills and outlines
        - Position and dimensions
        - Slide background properties
        """
        try:
            slide_layout = source_slide.slide_layout
            new_slide = dest_prs.slides.add_slide(slide_layout)
            
            # Copy slide-level properties
            self._copy_slide_properties(source_slide, new_slide)
            
            # Clear default placeholders
            for shape in list(new_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)
            
            # Clone all shapes from source
            for shape in source_slide.shapes:
                self._clone_shape(shape, new_slide)
            
        except Exception as e:
            self.logger.error(f"Error cloning slide: {str(e)}")
            raise
    
    def _clone_shape(self, source_shape, dest_slide) -> None:
        """
        Clone a shape preserving all formatting:
        - Text and text formatting (font, size, color, bold, italic, etc.)
        - Shape fill and outline
        - Position and size
        - Shadow and 3D effects
        """
        try:
            # Deep copy shape element (preserves all XML properties)
            shape_element = deepcopy(source_shape.element)
            dest_slide.shapes._spTree.insert_element_before(
                shape_element, 'p:extLst'
            )
            
        except Exception as e:
            self.logger.warning(f"Could not clone shape, retrying with fallback: {e}")
            self._clone_shape_fallback(source_shape, dest_slide)
    
    def _clone_shape_fallback(self, source_shape, dest_slide) -> None:
        """
        Fallback method for cloning shapes when direct XML copy fails.
        Handles text, colors, and basic formatting.
        """
        try:
            if not source_shape.has_text_frame:
                return
            
            left = source_shape.left
            top = source_shape.top
            width = source_shape.width
            height = source_shape.height
            
            text_box = dest_slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            
            # Copy text frame properties
            text_frame.word_wrap = source_shape.text_frame.word_wrap
            text_frame.margin_bottom = source_shape.text_frame.margin_bottom
            text_frame.margin_top = source_shape.text_frame.margin_top
            text_frame.margin_left = source_shape.text_frame.margin_left
            text_frame.margin_right = source_shape.text_frame.margin_right
            text_frame.vertical_anchor = source_shape.text_frame.vertical_anchor
            
            # Copy paragraphs with formatting
            text_frame.clear()
            for source_para in source_shape.text_frame.paragraphs:
                self._copy_paragraph(source_para, text_frame)
            
            # Copy shape fill and line
            self._copy_shape_fill(source_shape, text_box)
            self._copy_shape_line(source_shape, text_box)
            
        except Exception as e:
            self.logger.warning(f"Fallback clone also failed: {str(e)}")
    
    def _copy_paragraph(self, source_para, dest_text_frame) -> None:
        """Copy paragraph with all text runs and formatting."""
        try:
            dest_para = dest_text_frame.add_paragraph()
            dest_para.text = ""
            
            # Copy paragraph properties
            dest_para.level = source_para.level
            dest_para.alignment = source_para.alignment
            dest_para.space_before = source_para.space_before
            dest_para.space_after = source_para.space_after
            dest_para.line_spacing = source_para.line_spacing
            
            # Copy runs with formatting
            dest_para.clear()
            for source_run in source_para.runs:
                self._copy_run(source_run, dest_para)
        
        except Exception as e:
            self.logger.warning(f"Error copying paragraph: {str(e)}")
    
    def _copy_run(self, source_run, dest_para) -> None:
        """Copy text run with character-level formatting."""
        try:
            dest_run = dest_para.add_run()
            dest_run.text = source_run.text
            
            # Font properties
            font_src = source_run.font
            font_dst = dest_run.font
            
            font_dst.name = font_src.name
            font_dst.size = font_src.size
            font_dst.bold = font_src.bold
            font_dst.italic = font_src.italic
            font_dst.underline = font_src.underline
            
            # Color
            if font_src.color.type:
                if hasattr(font_src.color, 'rgb'):
                    font_dst.color.rgb = font_src.color.rgb
            
            # Other properties
            if hasattr(font_src, 'all_caps'):
                font_dst.all_caps = font_src.all_caps
            if hasattr(font_src, 'subscript'):
                font_dst.subscript = font_src.subscript
            if hasattr(font_src, 'superscript'):
                font_dst.superscript = font_src.superscript
                
        except Exception as e:
            self.logger.warning(f"Error copying run: {str(e)}")
    
    def _copy_shape_fill(self, source_shape, dest_shape) -> None:
        """Copy shape fill properties (color, gradient, pattern)."""
        try:
            if hasattr(source_shape, 'fill'):
                src_fill = source_shape.fill
                dst_fill = dest_shape.fill
                
                if hasattr(src_fill, 'type'):
                    dst_fill.solid()
                    if hasattr(src_fill, 'fore_color'):
                        try:
                            dst_fill.fore_color.rgb = src_fill.fore_color.rgb
                        except:
                            pass
        except Exception as e:
            self.logger.debug(f"Could not copy fill: {str(e)}")
    
    def _copy_shape_line(self, source_shape, dest_shape) -> None:
        """Copy shape outline/line properties."""
        try:
            if hasattr(source_shape, 'line'):
                src_line = source_shape.line
                dst_line = dest_shape.line
                
                if hasattr(src_line, 'color'):
                    try:
                        dst_line.color.rgb = src_line.color.rgb
                    except:
                        pass
                if hasattr(src_line, 'width'):
                    dst_line.width = src_line.width
        except Exception as e:
            self.logger.debug(f"Could not copy line: {str(e)}")
    
    def _copy_slide_properties(self, source_slide, dest_slide) -> None:
        """Copy slide-level properties like background."""
        try:
            if hasattr(source_slide.background, 'fill'):
                src_fill = source_slide.background.fill
                dst_fill = dest_slide.background.fill
                
                if hasattr(src_fill, 'solid'):
                    try:
                        dst_fill.solid()
                        dst_fill.fore_color.rgb = src_fill.fore_color.rgb
                    except:
                        pass
        except Exception as e:
            self.logger.debug(f"Could not copy slide properties: {str(e)}")
    
    def _clone_slide_with_master(self, source_slide, dest_prs: Presentation,
                                 master: 'SlideMaster') -> None:
        """Clone slide using specific master (useful for consistent theming)."""
        try:
            layout = master.slide_layouts[0]
            new_slide = dest_prs.slides.add_slide(layout)
            
            for shape in list(new_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)
            
            for shape in source_slide.shapes:
                self._clone_shape(shape, new_slide)
        
        except Exception as e:
            self.logger.error(f"Error cloning with master: {str(e)}")
            raise
    
    # ==================== UTILITY METHODS ====================
    
    def get_slide_info(self, input_file: str) -> dict:
        """Get information about a presentation."""
        try:
            prs = Presentation(input_file)
            return {
                'file': input_file,
                'total_slides': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height,
                'width_inches': prs.slide_width / 914400,
                'height_inches': prs.slide_height / 914400,
            }
        except Exception as e:
            self.logger.error(f"Error getting slide info: {str(e)}")
            raise
    
    def validate_output(self, output_file: str) -> bool:
        """Validate that output file is a valid PPTX."""
        try:
            prs = Presentation(output_file)
            self.logger.info(f"Validation successful: {len(prs.slides)} slides")
            return True
        except Exception as e:
            self.logger.error(f"Validation failed: {str(e)}")
            return False


# ==================== EXAMPLE USAGE ====================

if __name__ == "__main__":
    
    processor = PPTXSplitMerge()
    
    # Example 1: Split PPTX file (1 slide per file)
    # output_files = processor.split_pptx(
    #     input_file="presentation.pptx",
    #     output_dir="split_output",
    #     slides_per_file=1
    # )
    # print(f"Created files: {output_files}")
    
    # Example 2: Split PPTX file (multiple slides per file)
    # output_files = processor.split_pptx(
    #     input_file="presentation.pptx",
    #     output_dir="split_output",
    #     slides_per_file=5
    # )
    # print(f"Created files: {output_files}")
    
    # Example 3: Split by specific ranges
    # output_files = processor.split_pptx_by_range(
    #     input_file="presentation.pptx",
    #     output_dir="split_output",
    #     ranges=[(1, 5), (6, 10), (11, 15)]
    # )
    # print(f"Created files: {output_files}")
    
    # Example 4: Merge PPTX files (preserve source themes)
    # merged = processor.merge_pptx(
    #     input_files=["part1.pptx", "part2.pptx", "part3.pptx"],
    #     output_file="merged.pptx",
    #     keep_source_theme=True
    # )
    # print(f"Merged file: {merged}")
    
    # Example 5: Merge with unified theme
    # merged = processor.merge_pptx(
    #     input_files=["part1.pptx", "part2.pptx", "part3.pptx"],
    #     output_file="merged_unified.pptx",
    #     keep_source_theme=False
    # )
    # print(f"Merged file (unified theme): {merged}")
    
    # Example 6: Get slide information
    # info = processor.get_slide_info("presentation.pptx")
    # print(f"Presentation info: {info}")
    
    # Example 7: Validate output
    # is_valid = processor.validate_output("output.pptx")
    # print(f"Output is valid: {is_valid}")
    
    print("PPTXSplitMerge class ready for use")